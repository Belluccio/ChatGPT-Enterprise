import os
import openai
import tkinter as tk
from tkinter import filedialog, scrolledtext, END, Frame, Label
from dotenv import load_dotenv
import pdfplumber
from docx import Document
import openpyxl
import threading
import time
import ttkbootstrap as ttk
from ttkbootstrap.constants import PRIMARY, OUTLINE
from PIL import Image
import pytesseract
from pdf2image import convert_from_path
from pptx import Presentation
import numpy as np
import pandas as pd

# Asegúrate de instalar tiktoken
try:
    import tiktoken
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "tiktoken"])
    import tiktoken

from sklearn.metrics.pairwise import cosine_distances

# Configuración de la ruta de Tesseract (cambiar según tu sistema operativo)
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Cargar variables de entorno desde el archivo .env
load_dotenv('.env')

# Obtener la clave de API desde las variables de entorno
api_key = os.getenv('OPENAI_API_KEY')

if not api_key:
    raise ValueError("API key not found. Please check your .env file.")

# Configurar la clave API
openai.api_key = api_key

# Crear la ventana de la interfaz de usuario
root = ttk.Window(themename='litera')
root.title("Interfaz ChatGPT")
root.geometry("700x750")
root.configure(bg='#FFFFFF')

# Iconos Unicode
user_icon = "👤"
bot_icon = "💬"

# Variables para almacenar el contenido y embeddings de los documentos
document_content = ""
document_embeddings = pd.DataFrame(columns=['text', 'embedding'])
messages = [{"role": "system", "content": "Eres un asistente útil y responderás en español. Puedes procesar documentos de Word, PDF, Excel, PowerPoint, y extraer texto de imágenes."}]

# Función para obtener la respuesta de ChatGPT
def get_response(user_input):
    global messages
    try:
        messages.append({"role": "user", "content": user_input})
        context = create_context(user_input)
        response = openai.ChatCompletion.create(
            model="gpt-4-turbo",
            messages=messages + [{"role": "system", "content": f"Context: {context}"}],
            max_tokens=500,
            temperature=0.7
        )
        message_content = response['choices'][0]['message']['content'].strip()
        messages.append({"role": "assistant", "content": message_content})
        return message_content
    except openai.OpenAIError as e:
        return f"Error al llamar a la API de OpenAI: {e}"

# Función para crear el contexto basado en los embeddings
def create_context(question, max_len=1800):
    global document_embeddings
    if document_embeddings.empty:
        return "No se ha cargado ningún documento."

    question_embedding = openai.Embedding.create(input=question, model="text-embedding-ada-002")['data'][0]['embedding']
    document_embeddings['distance'] = distances_from_embeddings(np.array(question_embedding), document_embeddings['embedding'].values)

    context_texts = []
    cur_len = 0
    for _, row in document_embeddings.sort_values('distance').iterrows():
        cur_len += len(row['text'])
        if cur_len > max_len:
            break
        context_texts.append(row['text'])
    return "\n\n".join(context_texts)

# Función para extraer texto de un archivo PDF usando pdfplumber
def extract_text_from_pdf(file_path):
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                extracted_text = page.extract_text()
                if extracted_text:
                    text += extracted_text
                else:
                    text += ocr_pdf_page(page)
        return text
    except Exception as e:
        return f"Error al extraer texto del PDF: {e}"

# Función para realizar OCR en una página PDF que es una imagen
def ocr_pdf_page(page):
    try:
        image = page.to_image(resolution=300).original
        text = pytesseract.image_to_string(image, lang='spa')
        return text
    except Exception as e:
        return f"Error al realizar OCR en la página del PDF: {e}"

# Función para extraer texto de un archivo de Word
def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        text = ''
        for para in doc.paragraphs:
            text += para.text + '\n'
        return text
    except Exception as e:
        return f"Error al extraer texto del archivo Word: {e}"

# Función para extraer texto de un archivo de Excel
def extract_text_from_xlsx(file_path):
    try:
        wb = openpyxl.load_workbook(file_path)
        sheet = wb.active
        text = ''
        for row in sheet.iter_rows(values_only=True):
            text += ' '.join([str(cell) if cell is not None else '' for cell in row]) + '\n'
        return text
    except Exception as e:
        return f"Error al extraer texto del archivo Excel: {e}"

# Función para extraer texto de un archivo PowerPoint
def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        return text
    except Exception as e:
        return f"Error al extraer texto del archivo PowerPoint: {e}"

# Función para extraer texto de una imagen
def extract_text_from_image(file_path):
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image, lang='spa')
        return text
    except Exception as e:
        return f"Error al extraer texto de la imagen: {e}"

# Función para crear embeddings del texto
def create_embeddings(text):
    embeddings = []
    max_tokens = 500
    sentences = text.split('. ')
    n_tokens = [len(tiktoken.get_encoding("cl100k_base").encode(" " + sentence)) for sentence in sentences]

    tokens_so_far = 0
    chunk = []

    for sentence, token in zip(sentences, n_tokens):
        if tokens_so_far + token > max_tokens:
            chunk_text = ". ".join(chunk) + "."
            embedding = openai.Embedding.create(input=chunk_text, model="text-embedding-ada-002")['data'][0]['embedding']
            embeddings.append((chunk_text, embedding))
            chunk = []
            tokens_so_far = 0
        if token > max_tokens:
            continue
        chunk.append(sentence)
        tokens_so_far += token + 1

    if chunk:
        chunk_text = ". ".join(chunk) + "."
        embedding = openai.Embedding.create(input=chunk_text, model="text-embedding-ada-002")['data'][0]['embedding']
        embeddings.append((chunk_text, embedding))

    return embeddings

# Función para cargar y procesar archivos
def handle_file_upload():
    global document_content, document_embeddings
    file_path = filedialog.askopenfilename(filetypes=[("Todos los archivos", "*.pdf;*.docx;*.xlsx;*.pptx;*.jpg;*.png"), 
                                                      ("PDF files", "*.pdf"), 
                                                      ("Word files", "*.docx"), 
                                                      ("Excel files", "*.xlsx"), 
                                                      ("PowerPoint files", "*.pptx"), 
                                                      ("Image files", "*.jpg;*.png")])
    if not file_path:
        return

    add_message_to_chat_box("Sistema", "Cargando documento...", "center", "blue", "white", bot_icon)

    def process_file():
        nonlocal file_path
        if file_path.endswith('.pdf'):
            additional_content = extract_text_from_pdf(file_path)
        elif file_path.endswith('.docx'):
            additional_content = extract_text_from_docx(file_path)
        elif file_path.endswith('.xlsx'):
            additional_content = extract_text_from_xlsx(file_path)
        elif file_path.endswith('.pptx'):
            additional_content = extract_text_from_pptx(file_path)
        elif file_path.endswith('.jpg') or file_path.endswith('.png'):
            additional_content = extract_text_from_image(file_path)
        else:
            additional_content = f"Formato de archivo no soportado: {file_path}"

        if additional_content.startswith("Error"):
            add_message_to_chat_box("Sistema", additional_content, "center", "red", "white", bot_icon)
        else:
            document_content = additional_content
            embeddings = create_embeddings(document_content)
            document_embeddings_df = pd.DataFrame(embeddings, columns=['text', 'embedding'])
            document_embeddings_df['embedding'] = document_embeddings_df['embedding'].apply(np.array)
            document_embeddings = document_embeddings_df
            add_message_to_chat_box("Sistema", "Documento cargado exitosamente. Puedes hacer preguntas sobre el documento.", "center", "blue", "white", bot_icon)

    threading.Thread(target=process_file).start()

# Función para mostrar el efecto de "escribiendo..."
def show_typing_effect():
    typing_label.config(text="ChatGPT está escribiendo...")
    root.update_idletasks()

# Función para enviar el mensaje y obtener la respuesta
def send_message(event=None):
    user_input = input_box.get("1.0", END).strip()
    if user_input.lower() in ["salir", "exit", "quit"]:
        print("¡Adiós!")
        root.quit()
    else:
        add_message_to_chat_box("Tú", user_input, "right", "#e1ffc7", "black", user_icon)
        input_box.delete("1.0", END)
        threading.Thread(target=show_typing_effect).start()
        
        # Mostrar spinner
        spinner.pack(pady=5)
        spinner.start()
        
        # Obtener respuesta en un hilo separado
        threading.Thread(target=process_response, args=(user_input,)).start()

def process_response(user_input):
    response = get_response(user_input)
    
    # Detener y ocultar el spinner
    spinner.stop()
    spinner.pack_forget()
    
    # Ocultar la etiqueta de "escribiendo..."
    typing_label.config(text="")
    
    # Mostrar respuesta con efecto de escritura
    display_text(response, bot_icon)

def display_text(text, icon):
    frame = Frame(chat_box, bg="#d1e7dd", padx=10, pady=5)
    Label(frame, text=icon, bg="#d1e7dd", font=("Arial", 14)).pack(side=tk.LEFT, padx=5)
    text_label = Label(frame, text="", wraplength=500, bg="#d1e7dd", justify=tk.LEFT, anchor="w", font=("Arial", 12), padx=2, pady=5)
    text_label.pack(fill=tk.BOTH, expand=True)
    chat_box.window_create(END, window=frame)
    chat_box.insert(END, "\n")
    chat_box.window_create(END, window=tk.Frame(chat_box, height=1, bd=1, relief=tk.SUNKEN, bg="#f0f0f0", width=500))
    chat_box.insert(END, "\n")
    chat_box.yview(END)
    
    for char in text:
        text_label["text"] += char
        chat_box.update_idletasks()
        time.sleep(0.02)

# Función para añadir mensajes al chat con formato diferenciado
def add_message_to_chat_box(sender, message, align, bg_color, fg_color, icon):
    chat_box.config(state='normal')
    frame = Frame(chat_box, bg=bg_color, padx=10, pady=5)
    if align == "right":
        Label(frame, text=icon, bg=bg_color, font=("Arial", 14)).pack(side=tk.RIGHT, padx=5)
        text_label = Label(frame, text=message, wraplength=500, bg=bg_color, fg=fg_color, justify=tk.RIGHT, anchor="e", font=("Arial", 12), padx=2, pady=5)
    else:
        Label(frame, text=icon, bg=bg_color, font=("Arial", 14)).pack(side=tk.LEFT, padx=5)
        text_label = Label(frame, text=message, wraplength=500, bg=bg_color, fg=fg_color, justify=tk.LEFT, anchor="w", font=("Arial", 12), padx=2, pady=5)
    text_label.pack(fill=tk.BOTH, expand=True)
    chat_box.window_create(END, window=frame)
    chat_box.insert(END, "\n")
    chat_box.window_create(END, window=tk.Frame(chat_box, height=1, bd=1, relief=tk.SUNKEN, bg="#f0f0f0", width=500))
    chat_box.insert(END, "\n")
    chat_box.config(state='disabled')
    chat_box.yview(END)

# Crear un cuadro de texto para la conversación
chat_box = scrolledtext.ScrolledText(root, wrap=tk.WORD, state='disabled', width=80, height=20, bg='#FFFFFF', fg='#000000', font=("Arial", 12), spacing1=5)
chat_box.pack(padx=10, pady=10, fill=tk.BOTH, expand=True)

# Crear un cuadro de entrada para el usuario
input_frame = tk.Frame(root, bg='#FFFFFF')
input_frame.pack(padx=10, pady=10, fill=tk.X)

# Crear un botón de carga con un icono de texto Unicode (simulando FontAwesome)
upload_button = ttk.Button(input_frame, text="📁 Cargar", command=handle_file_upload, bootstyle=OUTLINE)
upload_button.pack(side=tk.LEFT, padx=(0, 10))

# Etiqueta informativa debajo del botón de carga
info_label = tk.Label(root, text="Acepta Documentos: Word, PDF, Excel, PowerPoint, Imágenes (JPG, PNG)", bg="#FFFFFF", fg="#000000", font=("Arial", 10))
info_label.pack(pady=(0, 20))

input_box = tk.Text(input_frame, height=2, width=50, bg='#f0f0f0', fg='#000000', font=("Arial", 12))
input_box.pack(side=tk.LEFT, padx=(10, 0), fill=tk.X, expand=True)
input_box.bind("<Return>", send_message)

# Crear un botón de enviar con bordes redondeados usando ttkbootstrap
send_button = ttk.Button(input_frame, text="Enviar", bootstyle=PRIMARY, command=send_message)
send_button.pack(side=tk.RIGHT, padx=(10, 0))

# Etiqueta para mostrar el efecto de "escribiendo..."
typing_label = tk.Label(root, text="", bg="#FFFFFF", fg="#000000", font=("Arial", 12))
typing_label.pack(pady=5)

# Crear un spinner para mostrar mientras se procesa la respuesta
spinner = ttk.Progressbar(root, mode='indeterminate', bootstyle=PRIMARY)
spinner.pack_forget()

# Función para cargar documentos desde una carpeta
def load_documents_from_folder(folder_path):
    global document_content, document_embeddings
    document_content = ""
    document_embeddings = pd.DataFrame(columns=['text', 'embedding'])
    
    for root_dir, dirs, files in os.walk(folder_path):
        for file in files:
            file_path = os.path.join(root_dir, file)
            if file.endswith('.pdf'):
                additional_content = extract_text_from_pdf(file_path)
            elif file.endswith('.docx'):
                additional_content = extract_text_from_docx(file_path)
            elif file.endswith('.xlsx'):
                additional_content = extract_text_from_xlsx(file_path)
            elif file.endswith('.pptx'):
                additional_content = extract_text_from_pptx(file_path)
            elif file.endswith('.jpg') or file.endswith('.png'):
                additional_content = extract_text_from_image(file_path)
            else:
                additional_content = f"Formato de archivo no soportado: {file_path}"

            if not additional_content.startswith("Error"):
                document_content += additional_content
                embeddings = create_embeddings(additional_content)
                embeddings_df = pd.DataFrame(embeddings, columns=['text', 'embedding'])
                embeddings_df['embedding'] = embeddings_df['embedding'].apply(np.array)
                document_embeddings = pd.concat([document_embeddings, embeddings_df], ignore_index=True)

    add_message_to_chat_box("Sistema", "Documentos cargados desde la carpeta inicial.", "center", "blue", "white", bot_icon)

# Cargar documentos de una carpeta al iniciar
document_folder = 'ruta/a/tu/carpeta'  # Reemplaza esto con la ruta a tu carpeta de documentos
load_documents_from_folder(document_folder)

root.mainloop()
